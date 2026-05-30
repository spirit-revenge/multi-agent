"""
Unified document processor for PDF, PPTX, and DOCX files.
Extracts text (semantically chunked), images, and tables (as Markdown).
"""

import io
import re
import logging
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime

from PIL import Image

logger = logging.getLogger(__name__)


def _safe_filename(name: str) -> str:
    """Replace URL-unsafe and Markdown-special characters with underscores."""
    # Replace characters unsafe in URLs or Markdown syntax
    safe = re.sub(r'[^\w\-.]', '_', name)
    # Collapse multiple underscores
    safe = re.sub(r'_+', '_', safe)
    # Strip leading/trailing underscores
    return safe.strip('_')


# ============================================================================
# Semantic chunking
# ============================================================================

def _semantic_chunk(text: str, min_size: int = 100, max_size: int = 1200) -> List[str]:
    """Split text into semantic chunks respecting paragraph and heading boundaries.

    Rules:
    - Prefer splitting at paragraph boundaries (double newline).
    - Treat markdown headings (##, ###) as hard boundaries.
    - Merge sub-paragraph chunks below min_size.
    - Force-split oversize chunks at sentence boundaries (。！？.!?).
    """
    if not text or not text.strip():
        return []

    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')

    # Split into raw paragraphs (double newline or heading lines)
    raw_paragraphs = re.split(r'\n\s*\n', text)
    paragraphs = []
    for para in raw_paragraphs:
        para = para.strip()
        if not para:
            continue
        # Check if this paragraph starts with a heading marker
        is_heading = bool(re.match(r'^#{1,3}\s', para))
        # Split long paragraphs at sentence boundaries
        if len(para) > max_size and not is_heading:
            sub = _split_at_sentence_boundary(para, max_size)
            paragraphs.extend(sub)
        else:
            paragraphs.append(para)

    # Merge small consecutive paragraphs
    merged = []
    for para in paragraphs:
        is_heading = bool(re.match(r'^#{1,3}\s', para))
        if (
            not is_heading
            and merged
            and len(merged[-1]) < min_size
            and not re.match(r'^#{1,3}\s', merged[-1])
        ):
            merged[-1] = merged[-1] + '\n' + para
        else:
            merged.append(para)

    return merged


def _split_at_sentence_boundary(text: str, max_size: int) -> List[str]:
    """Split a long text at sentence boundaries (。！？．.!?) near max_size."""
    if len(text) <= max_size:
        return [text]

    # Sentence-ending characters (Chinese + English)
    sentence_end = re.compile(r'([。！？．\.!\?])')
    parts = sentence_end.split(text)

    chunks = []
    buffer = ''
    for i in range(0, len(parts) - 1, 2):
        seg = parts[i] + (parts[i + 1] if i + 1 < len(parts) else '')
        if len(buffer) + len(seg) <= max_size:
            buffer += seg
        else:
            if buffer:
                chunks.append(buffer.strip())
            buffer = seg
    if buffer:
        chunks.append(buffer.strip())

    # If still oversize (no sentence boundary found), hard-split at max_size
    result = []
    for chunk in chunks:
        if len(chunk) > max_size:
            for i in range(0, len(chunk), max_size):
                result.append(chunk[i:i + max_size].strip())
        else:
            result.append(chunk)

    return [c for c in result if c]


# ============================================================================
# PDF processing (PyMuPDF)
# ============================================================================

def _process_pdf(file_path: Path) -> Dict:
    """Extract text, images, and tables from a PDF using PyMuPDF + pdfplumber."""
    import fitz  # PyMuPDF
    import pdfplumber

    texts: List[str] = []
    images: List[Tuple[Image.Image, str]] = []
    tables: List[str] = []

    # --- PyMuPDF: text + images ---
    doc = fitz.open(str(file_path))
    page_texts = []
    for page_num, page in enumerate(doc):
        # Text
        page_text = page.get_text()
        if page_text.strip():
            page_texts.append(page_text)

        # Images
        for img_index, img_info in enumerate(page.get_images(full=True)):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            try:
                pil_img = Image.open(io.BytesIO(image_bytes))
                # Skip tiny images (icons, bullets)
                if pil_img.width < 50 or pil_img.height < 50:
                    continue
                img_filename = f"{_safe_filename(file_path.stem)}_p{page_num + 1}_img{img_index}.png"
                images.append((pil_img, img_filename))
            except Exception:
                logger.debug("Skipping unreadable image on page %d", page_num + 1)

    doc.close()

    # Combine page texts with natural breaks
    full_text = '\n\n'.join(page_texts)
    texts = _semantic_chunk(full_text)

    # --- pdfplumber: tables (better table detection) ---
    try:
        with pdfplumber.open(str(file_path)) as pdf:
            for page_num, page in enumerate(pdf.pages):
                for table in page.extract_tables():
                    if not table or not any(any(cell for cell in row) for row in table):
                        continue
                    md_table = _table_to_markdown(table)
                    if md_table:
                        tables.append(md_table)
    except Exception as e:
        logger.warning("pdfplumber table extraction failed for %s: %s", file_path.name, e)

    return {"texts": texts, "images": images, "tables": tables}


# ============================================================================
# PPTX processing (python-pptx)
# ============================================================================

def _process_pptx(file_path: Path) -> Dict:
    """Extract text, images, and tables from a PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    texts: List[str] = []
    images: List[Tuple[Image.Image, str]] = []
    tables: List[str] = []
    slide_texts: List[str] = []

    prs = Presentation(str(file_path))
    for slide_idx, slide in enumerate(prs.slides):
        slide_paragraphs = []

        for shape in slide.shapes:
            # --- Text ---
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        slide_paragraphs.append(para_text)

            # --- Tables ---
            if shape.has_table:
                tbl = shape.table
                md_rows = []
                for row_idx, row in enumerate(tbl.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    md_rows.append(cells)
                md_table = _table_to_markdown(md_rows)
                if md_table:
                    tables.append(md_table)

            # --- Images ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    pil_img = Image.open(io.BytesIO(image_bytes))
                    if pil_img.width >= 50 and pil_img.height >= 50:
                        img_filename = f"{_safe_filename(file_path.stem)}_s{slide_idx + 1}_img{len(images)}.png"
                        images.append((pil_img, img_filename))
                except Exception:
                    pass

            # --- Group shapes (recursive) ---
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        if hasattr(child, 'has_text_frame') and child.has_text_frame:
                            for para in child.text_frame.paragraphs:
                                para_text = para.text.strip()
                                if para_text:
                                    slide_paragraphs.append(para_text)
                except Exception:
                    pass

        if slide_paragraphs:
            slide_texts.append('\n'.join(slide_paragraphs))

    full_text = '\n\n'.join(slide_texts)
    texts = _semantic_chunk(full_text)

    return {"texts": texts, "images": images, "tables": tables}


# ============================================================================
# DOCX processing (python-docx)
# ============================================================================

def _process_docx(file_path: Path) -> Dict:
    """Extract text, images, and tables from a DOCX file."""
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    texts: List[str] = []
    images: List[Tuple[Image.Image, str]] = []
    tables: List[str] = []
    paragraphs: List[str] = []

    doc = DocxDocument(str(file_path))

    # --- Text from paragraphs ---
    for para in doc.paragraphs:
        para_text = para.text.strip()
        if para_text:
            # Detect heading style
            if para.style and para.style.name and 'Heading' in para.style.name:
                level = para.style.name.replace('Heading ', '').strip()
                if level.isdigit():
                    level_num = int(level)
                    if level_num <= 3:
                        paragraphs.append(f"{'#' * level_num} {para_text}")
                        continue
            paragraphs.append(para_text)

    # --- Tables ---
    for table in doc.tables:
        md_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            md_rows.append(cells)
        md_table = _table_to_markdown(md_rows)
        if md_table:
            tables.append(md_table)

    # --- Images from inline shapes ---
    for rel in doc.part.rels.values():
        if rel.reltype == RT.IMAGE:
            try:
                image_bytes = rel.target_part.blob
                pil_img = Image.open(io.BytesIO(image_bytes))
                if pil_img.width >= 50 and pil_img.height >= 50:
                    img_filename = f"{_safe_filename(file_path.stem)}_img{len(images)}.png"
                    images.append((pil_img, img_filename))
            except Exception:
                pass

    full_text = '\n\n'.join(paragraphs)
    texts = _semantic_chunk(full_text)

    return {"texts": texts, "images": images, "tables": tables}


# ============================================================================
# Table → Markdown
# ============================================================================

def _table_to_markdown(rows: List[List[str]]) -> Optional[str]:
    """Convert a 2D list of cells to a GitHub-Flavored Markdown table string."""
    if not rows or len(rows) < 1:
        return None
    # Determine column count from the longest row
    num_cols = max(len(row) for row in rows)

    def _cell(text: str) -> str:
        return text.replace('\n', ' ').replace('|', '\\|').strip()

    # Header row
    header = rows[0]
    header_cells = [_cell(header[i]) if i < len(header) else '' for i in range(num_cols)]
    # Pad header
    while len(header_cells) < num_cols:
        header_cells.append('')

    # Separator
    separator = ['---'] * num_cols

    # Body rows
    body = []
    for row in rows[1:]:
        cells = [_cell(row[i]) if i < len(row) else '' for i in range(num_cols)]
        while len(cells) < num_cols:
            cells.append('')
        body.append(cells)

    if not body:
        return None

    lines = [
        '| ' + ' | '.join(header_cells) + ' |',
        '| ' + ' | '.join(separator) + ' |',
    ]
    for row in body:
        lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(lines)


# ============================================================================
# Main entry point
# ============================================================================

SUPPORTED_EXTENSIONS = {'.pdf', '.pptx', '.docx'}


def process_document(file_path: str) -> Dict:
    """Process a single document and return structured content.

    Returns:
        {
            "texts": [str, ...],       # Semantically chunked text paragraphs
            "images": [(PIL.Image, filename), ...],
            "tables": [markdown_str, ...],
        }
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    logger.info("Processing %s ...", path.name)

    if ext == '.pdf':
        result = _process_pdf(path)
    elif ext == '.pptx':
        result = _process_pptx(path)
    elif ext == '.docx':
        result = _process_docx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    logger.info(
        "  → %d text chunks, %d images, %d tables",
        len(result["texts"]),
        len(result["images"]),
        len(result["tables"]),
    )

    return result
