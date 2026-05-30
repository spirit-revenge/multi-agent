import os
from pathlib import Path
from typing import Type
from crewai.tools import BaseTool
from pydantic import BaseModel, Field
import PyPDF2
from pptx import Presentation

class ReadLocalFilesSchema(BaseModel):
    folder_path: str = Field(..., description="Path to folder containing PDF/PPT files")

class ReadLocalLectureFilesTool(BaseTool):
    name: str = "Read Local Lecture Files"
    description: str = "Reads all PDF and PPTX files from a folder and returns extracted text."
    args_schema: Type[BaseModel] = ReadLocalFilesSchema

    def _run(self, folder_path: str) -> str:
        folder = Path(folder_path)
        if not folder.exists():
            return f"错误：文件夹 '{folder_path}' 不存在"

        results = []
        for file_path in folder.iterdir():
            if file_path.suffix.lower() == '.pdf':
                content = self._read_pdf(file_path)
                results.append(f"\n--- {file_path.name} ---\n{content}\n")
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                content = self._read_pptx(file_path)
                results.append(f"\n--- {file_path.name} ---\n{content}\n")
        if not results:
            return "文件夹中未找到 PDF 或 PPT 文件"
        return "\n".join(results)

    def _read_pdf(self, path: Path) -> str:
        text = ""
        try:
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text
        except Exception as e:
            text = f"[读取 PDF 错误：{e}]"
        return text

    def _read_pptx(self, path: Path) -> str:
        text = ""
        try:
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        except Exception as e:
            text = f"[读取 PPTX 错误：{e}]"
        return text